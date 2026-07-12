import asyncio
import csv
import hashlib
import os
import threading
from collections.abc import Coroutine
from concurrent.futures import Future
from datetime import UTC, datetime
from typing import Any, cast
from uuid import UUID

import docx
import pdfplumber
import pptx
import pytesseract
from qdrant_client.http.models import PointStruct
from sentence_transformers import SentenceTransformer
from sqlalchemy import select

from backend.celery_worker.celery_app import celery_app
from backend.core.database import async_session_factory
from backend.core.qdrant import COLLECTION_NAME, init_qdrant_collection, qdrant_client
from backend.models.document import Document
from backend.repositories.document import ChunkRepository, DocumentRepository
from rag.chunking import RecursiveTextSplitter


def calculate_sha256(file_path: str) -> str:
    """Calculates the SHA-256 checksum of a file on disk."""
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        while chunk := f.read(8192):
            sha256.update(chunk)
    return sha256.hexdigest()


def parse_pdf(file_path: str) -> tuple[str, dict[str, str]]:
    """Parses a PDF file extracting text.

    Executes Tesseract OCR fallback on scanned pages.
    """
    text_content = []
    metadata: dict[str, str] = {}

    with pdfplumber.open(file_path) as pdf:
        # Extract metadata
        if pdf.metadata:
            metadata = {
                "author": pdf.metadata.get("Author", ""),
                "created_date": pdf.metadata.get("CreationDate", ""),
                "title": pdf.metadata.get("Title", ""),
            }

        # Extract text page by page
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            # If page text is empty/extremely short, fallback to Tesseract OCR
            if not page_text or len(page_text.strip()) < 50:
                try:
                    # Render page to PIL image
                    img = page.to_image(resolution=150).original
                    page_text = pytesseract.image_to_string(img)
                except Exception as e:
                    page_text = f"[OCR Failed: {e}]"

            if page_text:
                text_content.append(f"--- Page {i + 1} ---\n{page_text}")

    return "\n\n".join(text_content), metadata


def parse_docx(file_path: str) -> tuple[str, dict[str, str]]:
    """Parses a DOCX file, extracting text from paragraphs and tables."""
    doc = docx.Document(file_path)
    text_content = []

    # Extract metadata
    metadata = {
        "author": doc.core_properties.author or "",
        "created_date": (
            str(doc.core_properties.created) if doc.core_properties.created else ""
        ),
        "title": doc.core_properties.title or "",
    }

    # Extract paragraph texts
    for p in doc.paragraphs:
        if p.text.strip():
            text_content.append(p.text)

    # Extract tables texts
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                text_content.append(" | ".join(row_text))

    return "\n\n".join(text_content), metadata


def parse_pptx(file_path: str) -> tuple[str, dict[str, str]]:
    """Parses a PPTX presentation, extracting text from slides."""
    prs = pptx.Presentation(file_path)
    text_content = []

    metadata = {
        "author": prs.core_properties.author or "",
        "created_date": (
            str(prs.core_properties.created) if prs.core_properties.created else ""
        ),
        "title": prs.core_properties.title or "",
    }

    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_text.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        text_content.append("\n".join(slide_text))

    return "\n\n".join(text_content), metadata


def parse_csv(file_path: str) -> tuple[str, dict[str, str]]:
    """Parses a CSV tabular file."""
    text_content = []
    metadata: dict[str, str] = {}

    with open(file_path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if row:
                text_content.append(f"Row {i + 1}: " + " | ".join(row))

    return "\n".join(text_content), metadata


def run_async(coro: Coroutine[Any, Any, Any]) -> Any:
    """Safely runs an async coroutine, even if called inside an active event loop."""
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        return asyncio.run(coro)
    else:
        future: Future[Any] = Future()

        def _run() -> None:
            try:
                res = asyncio.run(coro)
                future.set_result(res)
            except Exception as e:
                future.set_exception(e)

        t = threading.Thread(target=_run)
        t.start()
        t.join()
        return future.result()


_embed_model = None


def get_embedding_model() -> SentenceTransformer:
    """Lazy-loader for the sentence-transformers embedding model."""
    global _embed_model
    if _embed_model is None:
        _embed_model = SentenceTransformer("all-MiniLM-L6-v2")
    return cast(SentenceTransformer, _embed_model)


@celery_app.task(name="backend.celery_worker.tasks.process_document_task")  # type: ignore
def process_document_task(document_id_str: str) -> str:
    """Asynchronous Celery task that processes uploaded documents.

    Handles duplicate checks, text parsing, OCR translation,
    recursive splitting, and database indexing.
    """
    document_id = UUID(document_id_str)

    async def _async_process() -> str:
        async with async_session_factory() as session:
            doc_repo = DocumentRepository(session)
            chunk_repo = ChunkRepository(session)

            # 1. Fetch document from database
            doc = await doc_repo.get(document_id)
            if not doc:
                return f"Document {document_id_str} not found."

            if not os.path.exists(doc.file_path):
                doc.status = "failed"
                await doc_repo.update(doc, doc)
                await session.commit()
                return f"File not found on disk: {doc.file_path}"

            # 2. Check for duplicate upload via SHA-256 hash
            try:
                checksum = calculate_sha256(doc.file_path)

                # Query duplicate documents BEFORE modifying doc.checksum
                stmt = select(Document).where(
                    Document.checksum == checksum, Document.id != doc.id
                )
                result = await session.execute(stmt)
                duplicate = result.scalars().first()

                if duplicate:
                    doc.status = "duplicate"
                    await doc_repo.update(doc, doc)
                    await session.commit()
                    # Clean up duplicate file on disk
                    if os.path.exists(doc.file_path):
                        os.remove(doc.file_path)
                    return (
                        f"Duplicate document detected (ID: {duplicate.id}). "
                        "Ingestion bypassed."
                    )

                # Set checksum if it is unique
                doc.checksum = checksum
            except Exception as e:
                doc.status = "failed"
                await doc_repo.update(doc, doc)
                await session.commit()
                return f"SHA-256 calculation failed: {e}"

            # 3. Parse content based on filetype
            ext = os.path.splitext(doc.filename)[1].lower()
            text = ""
            file_metadata: dict[str, str] = {}

            try:
                if ext == ".pdf":
                    text, file_metadata = parse_pdf(doc.file_path)
                elif ext in [".docx", ".doc"]:
                    text, file_metadata = parse_docx(doc.file_path)
                elif ext in [".pptx", ".ppt"]:
                    text, file_metadata = parse_pptx(doc.file_path)
                elif ext in [".csv", ".txt"]:
                    text, file_metadata = parse_csv(doc.file_path)
                else:
                    # Text/raw fallback
                    with open(doc.file_path, encoding="utf-8", errors="ignore") as f:
                        text = f.read()

                if not text.strip():
                    raise ValueError("No text could be extracted from this document.")

            except Exception as e:
                doc.status = "failed"
                await doc_repo.update(doc, doc)
                await session.commit()
                return f"Parsing failed for document {doc.filename}: {e}"

            # 4. Chunk text recursively
            splitter = RecursiveTextSplitter(chunk_size=1000, chunk_overlap=200)
            chunks = splitter.split_text(text)

            # 5. Insert Chunks into database and Qdrant vector db
            try:
                # Ensure vector database collection is initialized
                init_qdrant_collection()

                model = get_embedding_model()
                # Batched encoding for optimized performance
                embeddings = model.encode(chunks)

                points = []
                for idx, chunk_text in enumerate(chunks):
                    # Compile chunk metadata
                    chunk_meta = {
                        "filename": doc.filename,
                        "filetype": ext,
                        "ingested_at": datetime.now(UTC).isoformat(),
                        **file_metadata,
                    }
                    db_chunk = await chunk_repo.create(
                        {
                            "document_id": doc.id,
                            "content": chunk_text,
                            "index": idx,
                            "metadata_json": chunk_meta,
                        }
                    )

                    # Build Qdrant vector payload
                    payload = {
                        "document_id": str(doc.id),
                        "chunk_id": str(db_chunk.id),
                        "user_id": str(doc.user_id),
                        "filename": doc.filename,
                        "content": chunk_text,
                        "index": idx,
                        "ingested_at": chunk_meta["ingested_at"],
                    }
                    points.append(
                        PointStruct(
                            id=str(db_chunk.id),
                            vector=embeddings[idx].tolist(),
                            payload=payload,
                        )
                    )

                # Bulk upsert points to Qdrant
                qdrant_client.upsert(
                    collection_name=COLLECTION_NAME,
                    points=points,
                )

                # Mark document processed
                doc.status = "processed"
                await doc_repo.update(doc, doc)
                await session.commit()

            except Exception as e:
                doc.status = "failed"
                await doc_repo.update(doc, doc)
                await session.commit()
                return f"Database/vector chunk storage failed: {e}"

            return (
                f"Ingested {doc.filename} successfully. Generated {len(chunks)} chunks."
            )

    return cast(str, run_async(_async_process()))
